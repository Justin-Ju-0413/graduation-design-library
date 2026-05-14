import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = r'C:\Users\16084\Documents\Graduation_Design_Library\10_Final_Defense\Final_10min_EN\Final_Defense_15min_EN_v20_chapter_openers_CityUHK.pptx'
DST = r'C:\Users\16084\Documents\Graduation_Design_Library\10_Final_Defense\Final_10min_EN\Final_Defense_15min_EN_v22_final.pptx'

# CityUHK brand colors
CITYU_RED = RGBColor(0xB1, 0x11, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_BG = RGBColor(0x0B, 0x0B, 0x0B)
CARD_BG = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT2 = RGBColor(0x8B, 0x1A, 0x2B)
ACCENT3 = RGBColor(0x6E, 0x22, 0x2E)
ACCENT4 = RGBColor(0x50, 0x1A, 0x24)

print("Opening presentation...")
prs = Presentation(SRC)
SW = prs.slide_width
SH = prs.slide_height
print(f"Slide size: {SW}x{SH} EMU, Existing slides: {len(prs.slides)}")


def add_text_box(slide, left, top, width, height, text, font_size=Pt(18),
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri Light'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tf


def add_multi_text(slide, left, top, width, height, paragraphs_data):
    """paragraphs_data: list of (text, font_size, bold, color, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in paragraphs_data:
        text, font_size, bold, color, align = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Calibri Light'
    return tf


def add_rect(slide, left, top, width, height, fill_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def dark_slide_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG


blank_layout = prs.slide_layouts[0]

# ============================================================
# Chapter Data (updated with correct slide ranges)
# ============================================================
chapters = [
    {'num': '01', 'title': 'Motivation and Objective',
     'subtitle': 'Why build a custom-instruction CNN accelerator on RISC-V?',
     'slides': 'Slides 4–5', 'color': CITYU_RED},
    {'num': '02', 'title': 'Architecture and Design',
     'subtitle': 'How is the accelerator built and connected to the CPU?',
     'slides': 'Slides 7–10', 'color': ACCENT2},
    {'num': '03', 'title': 'Verification Evidence',
     'subtitle': 'What evidence proves the design works on real hardware?',
     'slides': 'Slides 12–17', 'color': ACCENT3},
    {'num': '04', 'title': 'Results and Future Work',
     'subtitle': 'What was achieved and what remains to be done?',
     'slides': 'Slides 19–23', 'color': ACCENT4},
]

# ============================================================
# NEW SLIDES (added at end, reordered later)
# ============================================================

# --- Slide: Agenda ---
print("Creating Agenda...")
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_text_box(s, Inches(0.72), Inches(0.5), Inches(11), Inches(0.6),
             'Presentation Outline', font_size=Pt(32), bold=True, color=WHITE)
add_text_box(s, Inches(0.72), Inches(1.15), Inches(11), Inches(0.4),
             'A complete evidence path from motivation to FPGA board validation.',
             font_size=Pt(14), color=LIGHT_GRAY)
add_rect(s, Inches(0.72), Inches(1.65), Inches(3.0), Pt(3), fill_color=CITYU_RED)

card_w, card_h = Inches(5.6), Inches(2.25)
gap_x, gap_y = Inches(0.35), Inches(0.25)
sx, sy = Inches(0.72), Inches(2.0)

for i, ch in enumerate(chapters):
    cx = sx + (i % 2) * (card_w + gap_x)
    cy = sy + (i // 2) * (card_h + gap_y)
    add_rect(s, cx, cy, card_w, card_h, fill_color=CARD_BG)
    add_rect(s, cx, cy, Pt(5), card_h, fill_color=ch['color'])
    add_text_box(s, cx + Inches(0.3), cy + Inches(0.2), Inches(1), Inches(0.5),
                 ch['num'], font_size=Pt(36), bold=True, color=ch['color'])
    add_text_box(s, cx + Inches(0.3), cy + Inches(0.7), card_w - Inches(0.6), Inches(0.45),
                 ch['title'], font_size=Pt(20), bold=True, color=WHITE)
    add_text_box(s, cx + Inches(0.3), cy + Inches(1.15), card_w - Inches(0.6), Inches(0.55),
                 ch['subtitle'], font_size=Pt(11), color=LIGHT_GRAY)
    add_text_box(s, cx + Inches(0.3), cy + Inches(1.8), card_w - Inches(0.6), Inches(0.3),
                 ch['slides'], font_size=Pt(10), color=MUTED)
add_text_box(s, Inches(0.72), Inches(6.8), Inches(11), Inches(0.3),
             'RISC-V Custom-Instruction CNN Accelerator  |  FYP Defense  |  May 13, 2026',
             font_size=Pt(10), color=LIGHT_GRAY)
idx_agenda = len(prs.slides) - 1

# --- Slides: 4 Chapter Dividers ---
idx_dividers = []
for ch in chapters:
    s = prs.slides.add_slide(blank_layout)
    dark_slide_bg(s)
    add_text_box(s, Inches(0.3), Inches(1.0), Inches(12), Inches(3.2),
                 ch['num'], font_size=Pt(220), bold=True, color=RGBColor(0x35, 0x35, 0x35))
    add_rect(s, Inches(0.8), Inches(4.3), Inches(2.0), Pt(4), fill_color=CITYU_RED)
    add_text_box(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.8),
                 ch['title'], font_size=Pt(44), bold=True, color=WHITE)
    add_text_box(s, Inches(0.8), Inches(5.35), Inches(11), Inches(0.5),
                 ch['subtitle'], font_size=Pt(18), color=LIGHT_GRAY)
    add_text_box(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.3),
                 f'CHAPTER {ch["num"]}  |  {ch["slides"]}',
                 font_size=Pt(11), color=MUTED)
    idx_dividers.append(len(prs.slides) - 1)

# --- Slide: Ch01 Scope & Deliverables ---
print("Creating Ch01 Scope...")
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_text_box(s, Inches(0.72), Inches(0.4), Inches(6), Inches(0.25),
             'SCOPE', font_size=Pt(10), bold=True, color=CITYU_RED)
add_text_box(s, Inches(0.72), Inches(0.7), Inches(11), Inches(0.6),
             'Project Scope and Deliverables', font_size=Pt(28), bold=True, color=WHITE)
add_rect(s, Inches(0.72), Inches(1.35), Inches(2.0), Pt(3), fill_color=CITYU_RED)

# Left: In Scope
add_text_box(s, Inches(0.72), Inches(1.7), Inches(5.5), Inches(0.35),
             'IN SCOPE', font_size=Pt(14), bold=True, color=GREEN)
scope_items = [
    'Design a RISC-V custom-instruction CNN accelerator via NICE interface',
    'Implement 4×4 INT8 systolic PE array with 6 custom instructions',
    'Integrate into Hummingbird E203 SoC on Davinci Pro A7-100T FPGA',
    'Validate end-to-end: RTL sim → full-SoC sim → FPGA board test',
    'Deliver a repeatable open-source prototype with documented evidence',
]
add_multi_text(s, Inches(0.72), Inches(2.15), Inches(5.5), Inches(3.8),
               [(f'• {t}', Pt(13), False, LIGHT_GRAY, PP_ALIGN.LEFT) for t in scope_items])

# Right: Out of Scope
add_text_box(s, Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.35),
             'OUT OF SCOPE', font_size=Pt(14), bold=True, color=LIGHT_GRAY)
out_items = [
    'ASIC synthesis or physical design',
    'Multi-core or multi-accelerator scaling',
    'Full network acceleration (FC layers run in software)',
    'Comparison with GPU/NPU/TPU performance',
    'Production-grade software toolchain or compiler',
]
add_multi_text(s, Inches(6.8), Inches(2.15), Inches(5.5), Inches(3.8),
               [(f'• {t}', Pt(13), False, RGBColor(0xBB, 0xBB, 0xBB), PP_ALIGN.LEFT) for t in out_items])

add_text_box(s, Inches(0.72), Inches(6.5), Inches(11), Inches(0.3),
             'Focused prototype validation, not an ASIC-scale accelerator claim.',
             font_size=Pt(12), color=RGBColor(0xAA, 0xAA, 0xAA))
idx_scope = len(prs.slides) - 1

# --- Slide: Ch02 Design Rationale ---
print("Creating Ch02 Design Rationale...")
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_text_box(s, Inches(0.72), Inches(0.4), Inches(6), Inches(0.25),
             'DESIGN RATIONALE', font_size=Pt(10), bold=True, color=CITYU_RED)
add_text_box(s, Inches(0.72), Inches(0.7), Inches(11), Inches(0.6),
             'Why These Design Choices?', font_size=Pt(28), bold=True, color=WHITE)
add_rect(s, Inches(0.72), Inches(1.35), Inches(2.0), Pt(3), fill_color=CITYU_RED)

decisions = [
    ('4×4 PE Array', 'INT8 systolic array',
     'Balances throughput vs. FPGA resource.\n4 PEs per row packs four INT8 into one 32-bit word.\nFits A7-100T with headroom (LUT 20.8%).'),
    ('Output-Stationary DF', 'Weight stays in PE',
     'Keeps partial sums local to each PE.\nMinimizes weight reload for small conv kernels.\nReduces readback overhead vs. weight-stationary.'),
    ('INT8 Precision', 'Edge inference standard',
     'Common in MCU-class edge AI (TFLite, CMSIS-NN).\nEnables 4× packing into 32-bit NICE registers.\nINT32 accumulation prevents overflow.'),
    ('NICE Interface', 'Instruction-level control',
     'No memory-mapped bus → no address decoding, no bus contention.\nDeterministic latency per instruction.\nDirect rs1/rs2 operand access from CPU regfile.'),
]

card_w2 = Inches(5.6)
card_h2 = Inches(2.1)
for i, (title, subtitle, body) in enumerate(decisions):
    col, row = i % 2, i // 2
    cx = Inches(0.72) + col * (card_w2 + Inches(0.35))
    cy = Inches(1.7) + row * (card_h2 + Inches(0.2))
    add_rect(s, cx, cy, card_w2, card_h2, fill_color=CARD_BG)
    add_rect(s, cx, cy, Pt(5), card_h2, fill_color=CITYU_RED)
    add_text_box(s, cx + Inches(0.3), cy + Inches(0.15), card_w2 - Inches(0.5), Inches(0.4),
                 title, font_size=Pt(16), bold=True, color=WHITE)
    add_text_box(s, cx + Inches(0.3), cy + Inches(0.55), card_w2 - Inches(0.5), Inches(0.25),
                 subtitle, font_size=Pt(10), color=CITYU_RED)
    add_text_box(s, cx + Inches(0.3), cy + Inches(0.85), card_w2 - Inches(0.5), Inches(1.1),
                 body, font_size=Pt(12), color=LIGHT_GRAY)
idx_design_rat = len(prs.slides) - 1

# --- Slide: New Evidence Map (replaces old slide 7) ---
print("Creating new Evidence Map...")
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_text_box(s, Inches(0.72), Inches(0.4), Inches(6), Inches(0.25),
             'EVIDENCE MAP', font_size=Pt(10), bold=True, color=CITYU_RED)
add_text_box(s, Inches(0.72), Inches(0.7), Inches(11), Inches(0.6),
             'Claims Are Tied to Board Evidence', font_size=Pt(28), bold=True, color=WHITE)
add_rect(s, Inches(0.72), Inches(1.35), Inches(2.0), Pt(3), fill_color=CITYU_RED)
add_text_box(s, Inches(0.72), Inches(1.5), Inches(11), Inches(0.45),
             'Each conclusion is backed by a specific simulation, ILA capture, UART output, or Vivado report.',
             font_size=Pt(14), color=LIGHT_GRAY)

evidence = [
    ('SoC Integration', 'E203 + NICE path verified in design and simulation',
     'SoC architecture + full-SoC simulation'),
    ('FPGA Boot', 'CPU reaches ITCM and UART is active',
     'hello_e203 UART + ILA PC trace'),
    ('NICE Instruction Path', 'Custom instructions execute on board',
     'NICE test + ILA signal capture'),
    ('CNN Correctness', 'SW, HW, and expected outputs match',
     'UART: 12, 23, 0, 19'),
    ('Performance & Fit', 'Kernel speedup and resource headroom',
     '5.282× + Vivado utilization'),
]

# Table header
hdr_y = Inches(2.1)
col1_x = Inches(0.72)
col2_x = Inches(3.0)
col3_x = Inches(7.5)
col_w = [Inches(2.1), Inches(4.3), Inches(4.3)]

add_rect(s, col1_x, hdr_y, col_w[0], Inches(0.35), fill_color=RGBColor(0x25, 0x25, 0x25))
add_rect(s, col2_x, hdr_y, col_w[1], Inches(0.35), fill_color=RGBColor(0x25, 0x25, 0x25))
add_rect(s, col3_x, hdr_y, col_w[2], Inches(0.35), fill_color=RGBColor(0x25, 0x25, 0x25))
add_text_box(s, col1_x + Inches(0.15), hdr_y + Inches(0.02), col_w[0], Inches(0.3),
             'Claim', font_size=Pt(12), bold=True, color=WHITE)
add_text_box(s, col2_x + Inches(0.15), hdr_y + Inches(0.02), col_w[1], Inches(0.3),
             'Evidence', font_size=Pt(12), bold=True, color=WHITE)
add_text_box(s, col3_x + Inches(0.15), hdr_y + Inches(0.02), col_w[2], Inches(0.3),
             'Figure / Output', font_size=Pt(12), bold=True, color=WHITE)

# Table rows
for i, (claim, ev, fig) in enumerate(evidence):
    ry = hdr_y + Inches(0.35) + i * Inches(0.58)
    bg_c = RGBColor(0x13, 0x13, 0x13) if i % 2 == 0 else RGBColor(0x1A, 0x1A, 0x1A)
    add_rect(s, col1_x, ry, col_w[0], Inches(0.58), fill_color=bg_c)
    add_rect(s, col2_x, ry, col_w[1], Inches(0.58), fill_color=bg_c)
    add_rect(s, col3_x, ry, col_w[2], Inches(0.58), fill_color=bg_c)
    add_text_box(s, col1_x + Inches(0.15), ry + Inches(0.12), col_w[0] - Inches(0.2), Inches(0.35),
                 claim, font_size=Pt(14), bold=True, color=WHITE)
    add_text_box(s, col2_x + Inches(0.15), ry + Inches(0.12), col_w[1] - Inches(0.2), Inches(0.35),
                 ev, font_size=Pt(13), color=LIGHT_GRAY)
    add_text_box(s, col3_x + Inches(0.15), ry + Inches(0.12), col_w[2] - Inches(0.2), Inches(0.35),
                 fig, font_size=Pt(13), color=CITYU_RED)

add_text_box(s, Inches(0.72), Inches(6.5), Inches(11), Inches(0.3),
             'Staged engineering evidence — not one isolated screenshot.',
             font_size=Pt(13), color=LIGHT_GRAY)
idx_evidence_new = len(prs.slides) - 1

# --- Slide: Ch04 Results vs Design Targets ---
print("Creating Ch04 Results vs Targets...")
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_text_box(s, Inches(0.72), Inches(0.4), Inches(6), Inches(0.25),
             'RESULTS VS. TARGETS', font_size=Pt(10), bold=True, color=CITYU_RED)
add_text_box(s, Inches(0.72), Inches(0.7), Inches(11), Inches(0.6),
             'Design Targets: All Met or Exceeded', font_size=Pt(28), bold=True, color=WHITE)
add_rect(s, Inches(0.72), Inches(1.35), Inches(2.0), Pt(3), fill_color=CITYU_RED)

targets = [
    ('Convolution Speedup', '≥5×', '5.282×', True,
     '287 accelerator cycles vs. 1,516 CPU cycles'),
    ('Accuracy Loss', '< 5%', '~1.7%', True,
     '10/10 LeNet-5 test images correct on FPGA'),
    ('FPGA Resource Fit', 'Fit on A7-100T', 'LUT 20.8%\nFF 10.1%\nBRAM 26.3%', True,
     'Plenty of headroom for larger PE array'),
    ('Timing Closure', 'Meet 16 MHz', 'WNS = 12.468 ns\n@ 16 MHz', True,
     'ILA debug builds also converged at 50 MHz'),
    ('End-to-End Demo', 'FPGA runs CNN', 'LeNet-5 running\non FPGA', True,
     'Complete flow: firmware → NICE → accelerator → UART'),
]

row_h = Inches(0.85)
start_y = Inches(1.7)
cols_x = [Inches(0.72), Inches(3.5), Inches(5.5), Inches(7.0), Inches(7.8)]

# Header
for j, (hdr, cx, cw) in enumerate([
    ('Metric', cols_x[0], cols_x[1] - cols_x[0]),
    ('Target', cols_x[1], cols_x[2] - cols_x[1]),
    ('Achieved', cols_x[2], cols_x[3] - cols_x[2]),
    ('', cols_x[3], cols_x[4] - cols_x[3]),
]):
    add_rect(s, cx, start_y, cw, Inches(0.3), fill_color=RGBColor(0x25, 0x25, 0x25))
    if hdr:
        add_text_box(s, cx + Inches(0.1), start_y + Inches(0.02), cw, Inches(0.25),
                     hdr, font_size=Pt(11), bold=True, color=WHITE)

for i, (metric, target, achieved, met, note) in enumerate(targets):
    ry = start_y + Inches(0.3) + i * row_h
    bg_c = RGBColor(0x13, 0x13, 0x13) if i % 2 == 0 else RGBColor(0x1A, 0x1A, 0x1A)

    add_rect(s, cols_x[0], ry, cols_x[4] - cols_x[0], row_h, fill_color=bg_c)
    add_rect(s, cols_x[0], ry, Pt(4), row_h, fill_color=GREEN if met else CITYU_RED)

    add_text_box(s, cols_x[0] + Inches(0.2), ry + Inches(0.15), cols_x[1] - cols_x[0] - Inches(0.3), Inches(0.55),
                 metric, font_size=Pt(14), bold=True, color=WHITE)
    add_text_box(s, cols_x[1] + Inches(0.1), ry + Inches(0.15), cols_x[2] - cols_x[1] - Inches(0.2), Inches(0.55),
                 target, font_size=Pt(13), color=LIGHT_GRAY)
    add_text_box(s, cols_x[2] + Inches(0.1), ry + Inches(0.1), cols_x[3] - cols_x[2] - Inches(0.2), Inches(0.65),
                 achieved, font_size=Pt(15), bold=True, color=GREEN)
    add_text_box(s, cols_x[3] + Inches(0.2), ry + Inches(0.15), Inches(4.5), Inches(0.55),
                 note, font_size=Pt(11), color=MUTED)

idx_results_vs = len(prs.slides) - 1

# --- Slide: Thank You ---
print("Creating Thank You...")
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_text_box(s, Inches(0.0), Inches(1.8), Inches(13.333), Inches(1.5),
             'Thank You', font_size=Pt(60), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_rect(s, Inches(5.9), Inches(3.3), Inches(1.5), Pt(3), fill_color=CITYU_RED)
add_text_box(s, Inches(0.0), Inches(3.6), Inches(13.333), Inches(0.8),
             'Questions & Discussion', font_size=Pt(28), color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.0), Inches(5.2), Inches(13.333), Inches(0.8),
             'RISC-V Custom-Instruction CNN Accelerator\nJustin JU  |  May 13, 2026',
             font_size=Pt(14), color=MUTED, alignment=PP_ALIGN.CENTER)
idx_thanks = len(prs.slides) - 1

# ============================================================
# Reorder slides
# ============================================================
print("Reordering slides...")
# Old:  0=Cover, 1=OldAgenda(rm), 2=Motiv, 3=Arch, 4=Instr, 5=PE,
#       6=VerStrat, 7=OldEvidence(rm), 8=Impl, 9=FPGA, 10=NICE, 11=Correct,
#       12=CNNtest, 13=LeNet, 14=Perf, 15=Limit
# New:  16=Agenda, 17=Ch01Div, 18=Ch02Div, 19=Ch03Div, 20=Ch04Div,
#       21=Scope, 22=DesignRat, 23=NewEvidence, 24=ResultsVs, 25=Thanks

new_order_indices = [
    0,      # 1. Cover
    16,     # 2. NEW Agenda
    17,     # 3. NEW Ch01 Divider
    2,      # 4. Motivation (Ch01)
    21,     # 5. NEW Ch01 Scope
    18,     # 6. NEW Ch02 Divider
    3,      # 7. Architecture
    4,      # 8. Instructions
    5,      # 9. PE Array
    22,     # 10. NEW Ch02 Design Rationale
    19,     # 11. NEW Ch03 Divider
    6,      # 12. Verification Strategy
    23,     # 13. NEW Evidence Map
    8,      # 14. Implementation Flow
    9,      # 15. FPGA Validation
    10,     # 16. NICE Path
    11,     # 17. Correctness
    20,     # 18. NEW Ch04 Divider
    12,     # 19. CNN Board Test
    13,     # 20. LeNet-5
    14,     # 21. Performance
    15,     # 22. Limitations
    24,     # 23. NEW Results vs Targets
    25,     # 24. NEW Thank You
]

removed_indices = {1, 7}  # old agenda, old evidence map

pres_xml = prs.element
sldIdLst = pres_xml.find(qn('p:sldIdLst'))
sld_ids = list(sldIdLst.findall(qn('p:sldId')))

# Remove unwanted slides (process in reverse order)
for ri in sorted(removed_indices, reverse=True):
    sldIdLst.remove(sld_ids[ri])
    del sld_ids[ri]

def map_pos(old_idx):
    """Map old slide index to position in sld_ids after removals"""
    return old_idx - sum(1 for r in removed_indices if r < old_idx)

reordered = [sld_ids[map_pos(i)] for i in new_order_indices]

for sld in list(sldIdLst.findall(qn('p:sldId'))):
    sldIdLst.remove(sld)
for sld in reordered:
    sldIdLst.append(sld)

print(f"Reordered: {len(reordered)} slides (removed old agenda + old evidence map)")

# ============================================================
# Save
# ============================================================
print(f"Saving to {DST}...")
prs.save(DST)
print("Done! v22_final.pptx")
