"""Generate final-defense evidence index, English PPTX, and slide previews.

The deck is intentionally conservative and editable. It reuses thesis figures
that are generated from real project evidence.
"""
from __future__ import annotations

import json
import re
import zipfile
from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
THESIS = ROOT / "thesis_latex"
FIG = THESIS / "figures"
DEFENSE = ROOT / "10_Final_Defense"
PRES = ROOT / "05_Presentation" / "Final"
PREVIEW = ROOT / "05_Presentation" / "Final" / "previews"
EVIDENCE = DEFENSE / "Evidence_Package"

TITLE = "RISC-V Custom Instruction Based Lightweight CNN Accelerator"
SUBTITLE = "FPGA Prototype Validation on E203 + NICE"

INK = RGBColor(31, 41, 51)
MUTED = RGBColor(91, 103, 115)
BLUE = RGBColor(47, 111, 179)
GREEN = RGBColor(35, 131, 90)
AMBER = RGBColor(183, 121, 31)
RED = RGBColor(180, 35, 24)
LIGHT = RGBColor(248, 250, 252)
LINE = RGBColor(203, 213, 225)


def font(size, bold=False, color=INK):
    return {"size": Pt(size), "bold": bold, "color": color}


def set_run(run, size=24, bold=False, color=INK, name="Aptos"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=24, bold=False, color=INK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return shape


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.34, 12.0, 0.55, title, 27, True, INK)
    if subtitle:
        add_textbox(slide, 0.58, 0.91, 11.7, 0.32, subtitle, 10.5, False, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.28), Inches(1.0), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_footer(slide, n):
    add_textbox(slide, 0.55, 7.18, 8.5, 0.22, "Final Year Project Defense | Evidence-backed FPGA validation", 7.5, False, MUTED)
    add_textbox(slide, 12.2, 7.18, 0.65, 0.22, f"{n:02d}", 7.5, False, MUTED, PP_ALIGN.RIGHT)


def fit_image(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    nw, nh = iw * scale, ih * scale
    px = x + (w - nw) / 2
    py = y + (h - nh) / 2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(nw), Inches(nh))


def add_bullets(slide, x, y, w, h, bullets, size=17):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
    return shape


def add_note(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_card(slide, x, y, w, h, title, body, color=BLUE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT
    rect.line.color.rgb = LINE
    add_textbox(slide, x + 0.18, y + 0.12, w - 0.36, 0.26, title, 12, True, color)
    add_textbox(slide, x + 0.18, y + 0.47, w - 0.36, h - 0.55, body, 10.5, False, INK)


def add_source(slide, text):
    add_textbox(slide, 0.58, 6.85, 11.8, 0.22, text, 7.4, False, MUTED)


def build_deck():
    PRES.mkdir(parents=True, exist_ok=True)
    PREVIEW.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 cover
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(246, 248, 251)
    add_textbox(s, 0.78, 0.78, 9.4, 1.45, TITLE, 31, True, INK)
    add_textbox(s, 0.82, 2.36, 6.4, 0.46, SUBTITLE, 16, False, MUTED)
    fit_image(s, FIG / "fig3_1_soc_architecture.png", 5.9, 3.05, 6.55, 3.1)
    add_textbox(s, 0.82, 6.42, 6.2, 0.26, "JU JIAXING | BEng Computer Engineering | May 2026", 10, False, MUTED)
    add_note(s, "Opening: this project demonstrates a RISC-V custom-instruction CNN accelerator and validates it from RTL simulation to FPGA board execution.")

    slides = []
    slides.append(("Why this project matters", "Edge CNN inference needs more throughput than a small CPU can provide, but a full custom accelerator must still remain verifiable.", ["CNN convolution is MAC-heavy and expensive for a scalar CPU.", "RISC-V custom instructions let the accelerator stay close to the execution pipeline.", "The project target is not only speedup, but a reproducible FPGA validation chain."], None, "Set motivation: the design is about acceleration and evidence, not just RTL implementation."))
    slides.append(("System architecture", "The accelerator is attached through the NICE coprocessor interface rather than memory-mapped I/O.", [], "fig3_1_soc_architecture.png", "Explain that NICE carries operands and control directly, while AHB remains responsible for memory and peripherals."))
    slides.append(("Custom instruction interface", "Six custom0 instructions form a small software-visible ISA for convolution control.", ["CFG/CLEAR configure and reset the accelerator state.", "WLOAD/DLOAD serialize 4x4 INT8 data into the PE array.", "COMP/RSTAT trigger compute and read the accumulated result."], "fig3_2b_instruction_table.png", "Walk through the six instructions and emphasize the rs2 index semantics."))
    slides.append(("Accelerator datapath", "A 4x4 INT8 PE array exposes parallel MAC throughput through register-fed custom instructions.", [], "fig3_4_pe_array.png", "Explain WLOAD columns, DLOAD rows, output-stationary accumulation, and RSTAT readback."))
    slides.append(("Verification strategy", "The same evidence chain supports the thesis and the defense.", [], "fig3_7_verification_chain.png", "Emphasize the closed chain: RTL, full-SoC, hello board, CNN board, rs2 regression."))
    slides.append(("Board bring-up evidence", "The hello_e203 run closed the CPU boot path with UART and ILA evidence.", ["MROM boot redirects execution to ITCM.", "UART proves firmware-level execution.", "ILA confirms PC activity and board-level reset/clock health."], "fig4_1_ila_pc_trace.png", "Use this slide to explain how ILA made the board bring-up observable."))
    slides.append(("CNN board result", "The CNN v1 demo passed on the FPGA: hardware output matched software and expected values.", [], "fig_uart_output.png", "State the board result plainly: SW output, HW output, expected output all match; speedup is 5.282x."))
    slides.append(("Performance and fit", "The accelerator reaches the target convolution speedup while fitting comfortably in the A7-100T.", [], "fig4_3_speedup_bar.png", "Start with speedup, then mention resource headroom and why end-to-end inference is still limited by software FC layers."))
    slides.append(("NICE rs2 bug case study", "The key integration bug was not in the PE array: it was an E203 decoder optimization conflicting with NICE operand semantics.", ["Problem: rs2 index 0 encoded as x0, so the decoder skipped rs2 capture.", "Fix: for NICE instructions, use nice_need_rs2 directly before x0 gating.", "Regression: rebuilt cnn_sysclk_ila and re-ran CNN v1 board demo successfully."], None, "Explain this as the engineering insight: protocol integration details matter as much as datapath correctness."))
    slides.append(("Contributions and next work", "The project closes a reproducible custom-instruction accelerator prototype and leaves clear expansion paths.", ["Implemented a 4x4 INT8 CNN accelerator with six NICE instructions.", "Closed RTL, full-SoC, hello_e203 board, and CNN v1 board validation.", "Future work: FC acceleration, higher SoC clock, larger benchmark coverage."], None, "Close with what is achieved and what remains deliberately scoped as future work."))

    for idx, (title, subtitle, bullets, img, note) in enumerate(slides, start=2):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        add_title(s, title, subtitle)
        if img:
            if bullets:
                fit_image(s, FIG / img, 0.66, 1.65, 7.35, 4.85)
                add_bullets(s, 8.35, 1.75, 4.25, 3.85, bullets, 15)
            else:
                fit_image(s, FIG / img, 0.92, 1.55, 11.4, 5.1)
        else:
            if title == "NICE rs2 bug case study":
                add_card(s, 0.9, 1.65, 3.45, 3.75, "Root cause", "The E203 decoder treated rs2=x0 as no operand needed. NICE WLOAD/DLOAD use rs2 as a vector-bank index, so index 0 was incorrectly gated.", RED)
                add_card(s, 4.95, 1.65, 3.15, 3.75, "Fix", "For NICE instructions, rv32_need_rs2 follows nice_need_rs2 directly. Standard RISC-V x0 gating remains unchanged.", BLUE)
                add_card(s, 8.7, 1.65, 3.35, 3.75, "Board regression", "cnn_sysclk_ila rebuilt cleanly. UART confirmed HW/SW/Expected match and 5.282x speedup after the fix.", GREEN)
            elif title == "Why this project matters":
                add_bullets(s, 0.95, 1.75, 5.4, 3.8, bullets, 20)
                fit_image(s, FIG / "fig3_6_build_pipeline.png", 6.65, 1.68, 5.6, 3.9)
            else:
                add_bullets(s, 0.95, 1.75, 11.0, 3.8, bullets, 20)
        add_source(s, "Evidence: thesis_latex/figures + 04_Experiments board/simulation logs.")
        add_footer(s, idx)
        add_note(s, note)

    out = PRES / "FYP_Final_Defense_English_Draft.pptx"
    prs.save(out)
    remove_notes_slide_number_placeholders(out)
    return out


def remove_notes_slide_number_placeholders(pptx_path: Path):
    """Remove default note-page slide-number placeholders from generated PPTX."""
    with zipfile.ZipFile(pptx_path, "r") as zin:
        items = {name: zin.read(name) for name in zin.namelist()}

    for name, blob in list(items.items()):
        if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", name):
            text = blob.decode("utf-8")
            text = re.sub(
                r'<p:sp><p:nvSpPr><p:cNvPr[^>]*name="Slide Number Placeholder[\s\S]*?</p:sp>',
                "",
                text,
            )
            items[name] = text.encode("utf-8")

    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, blob in items.items():
            zout.writestr(name, blob)
    tmp.replace(pptx_path)


def evidence_items():
    return [
        ("RTL simulation", "NICE unit behavior verified", ROOT / "04_Experiments/RTL_Simulation/2026-04-23_baseline_rerun/rtl_sim_results.txt"),
        ("Full-SoC simulation", "SDK/full-SoC baseline closed", ROOT / "04_Experiments/FullSoC_Simulation/2026-04-23_baseline_rerun/fullsoc_sim_results.txt"),
        ("hello_e203 board", "UART and ILA boot-chain evidence", ROOT / "04_Experiments/Board_BringUp/2026-04-28_board_connection_check/hello_e203_board_artifacts/ila_summary.txt"),
        ("hello_e203 UART", "board program reached boot/uart ok/loop", ROOT / "04_Experiments/Board_BringUp/2026-04-28_board_connection_check/hello_e203_board_artifacts/uart_output.txt"),
        ("CNN v1 board", "HW/SW/Expected match and speedup", ROOT / "04_Experiments/Board_BringUp/2026-05-09_nice_rs2_fix_verification/uart_output.txt"),
        ("CNN ILA", "post-fix board ILA capture", ROOT / "04_Experiments/Board_BringUp/2026-05-09_nice_rs2_fix_verification/ila_capture.csv"),
        ("rs2 fix summary", "decoder fix board verification", ROOT / "04_Experiments/Board_BringUp/2026-05-09_nice_rs2_fix_verification/BOARD_VERIFICATION.md"),
        ("Baselines", "active repo commit pair", ROOT / "02_Source_Repos/CURRENT_BASELINES.md"),
    ]


def write_evidence_package():
    EVIDENCE.mkdir(parents=True, exist_ok=True)
    rows = []
    for claim, result, path in evidence_items():
        rows.append((claim, result, path))
    md = ["# Final Defense Evidence Index", "", "Every thesis/PPT claim below points to an existing evidence file.", "", "| Claim | Result | Evidence path |", "|---|---|---|"]
    for claim, result, path in rows:
        status = "exists" if path.exists() else "missing"
        rel = path.relative_to(ROOT) if path.exists() or str(path).startswith(str(ROOT)) else path
        md.append(f"| {claim} | {result} | `{rel}` ({status}) |")
    md.extend(
        [
            "",
            "## Core Reporting Line",
            "",
            "`RTL -> full-SoC -> hello_e203 board -> CNN/NICE board -> rs2 bug fix case study`",
            "",
            "## Notes",
            "",
            "- Figures in the thesis and deck are redrawn summaries; raw CSV/log/screenshot evidence remains in `04_Experiments/`.",
            "- Do not change numeric claims unless the corresponding evidence file is updated first.",
        ]
    )
    out = EVIDENCE / "EVIDENCE_INDEX.md"
    out.write_text("\n".join(md), encoding="utf-8")
    return out


def render_pptx_previews(pptx_path: Path):
    """Render preview PNGs using a lightweight PIL approximation.

    This does not claim PowerPoint parity; it gives a readable contact sheet
    for quick visual QA when Office/LibreOffice rendering is unavailable.
    """
    PREVIEW.mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx_path)
    font_title = ImageFont.truetype("arial.ttf", 36)
    font_body = ImageFont.truetype("arial.ttf", 20)
    paths = []
    for i, slide in enumerate(prs.slides, start=1):
        im = Image.new("RGB", (1600, 900), "white")
        draw = ImageDraw.Draw(im)
        draw.rectangle([0, 0, 1600, 900], fill=(248, 250, 252) if i == 1 else "white")
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        y = 55
        if texts:
            draw.text((65, y), texts[0][:95], font=font_title, fill=(31, 41, 51))
            y += 72
            for t in texts[1:8]:
                for line in t.splitlines():
                    for wrapped in wrap(line, 88):
                        draw.text((80, y), wrapped, font=font_body, fill=(70, 80, 92))
                        y += 28
                    y += 4
        draw.text((1450, 850), f"{i:02d}", font=font_body, fill=(91, 103, 115))
        out = PREVIEW / f"slide_{i:02d}.png"
        im.save(out)
        paths.append(out)
    return paths


def inspect_pptx(pptx_path: Path):
    with zipfile.ZipFile(pptx_path) as zf:
        names = zf.namelist()
        media_count = sum(1 for n in names if n.startswith("ppt/media/"))
        visible_xml = [
            n
            for n in names
            if re.match(r"ppt/(slides|notesSlides)/.*\.xml$", n)
        ]
        xml_text = "\n".join(
            zf.read(n).decode("utf-8", errors="ignore") for n in visible_xml
        )
    return {
        "slides": len(Presentation(pptx_path).slides),
        "notes": len([n for n in names if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", n)]),
        "media_count": media_count,
        "visible_has_slide_number_placeholder": "sldNum" in xml_text or "Slide Number" in xml_text,
        "visible_has_chinese": any("\u4e00" <= ch <= "\u9fff" for ch in xml_text),
        "visible_has_board_evidence_pending": "board evidence pending" in xml_text.lower(),
    }


def main():
    evidence = write_evidence_package()
    deck = build_deck()
    previews = render_pptx_previews(deck)
    report = inspect_pptx(deck)
    report["deck"] = str(deck)
    report["evidence_index"] = str(evidence)
    report["previews"] = [str(p) for p in previews]
    report_path = PRES / "FYP_Final_Defense_English_Draft_QA.json"
    report_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(json.dumps(report, indent=2))


if __name__ == "__main__":
    main()
